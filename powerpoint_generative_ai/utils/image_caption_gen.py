import torch
from PIL.Image import Image
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from transformers import BlipProcessor, BlipForConditionalGeneration
from typing import List

class ImageCaptionGenerator:
    def __init__(self, model: str = "Salesforce/blip-image-captioning-base", device: str = "cuda"):
        self.device = torch.device(device)
        self.processor = BlipProcessor.from_pretrained(model)
        self.model = BlipForConditionalGeneration.from_pretrained(model).to(self.device)

    def infer(self, image: Image, prompt: str = ""):
        inputs = self.processor(image.convert("RGB"), prompt, return_tensors="pt")
        for k, v in inputs.items():
            inputs[k] = v.to(self.device)

        out = self.model.generate(**inputs)
        caption = self.processor.decode(out[0], skip_special_tokens=True)
        return caption.lstrip(prompt).strip()


def find_image_shape(shape: BaseShape, path: List[int] = None) -> List[int]:
    """
    Recursively searches a shape object for an image, mainly used for GROUPs. Returns
    a path to the image as a list of integers
    """
    # Create a new path if it doesn't exist
    if path is None:
        path = []
    
    # Base case: if the shape_type is picture
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return path

    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
        if hasattr(shape, "image"):
            return path

    # If the shape is a group 
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # Check each sub shape
        for i, sub_shape in enumerate(shape.shapes):
            result = find_image_shape(sub_shape, path + [i])
            # If a sub shape has shape_type 13, return the path
            if result is not None:
                return result

    # If no picture is found in this branch, return None
    return None


def get_image_shape_from_path(shape: BaseShape, path: List[int]) -> BaseShape:
    """
    Uses the path from the `find_image_shape` function to search for the image and
    return it
    """
    for index in path:
        shape = shape.shapes[index]
    return shape