import json
import openai
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import List
from .domain.constants import MAX_CONTENT_LENGTH
from .domain.exceptions import InvalidModel
from .domain.prompts import (
    DECK_CREATION_SYSTEM_PROMPT,
    CHART_DATA_IDENTIFICATION,
    BEST_CHART_FOR_DATA_SYSTEM_PROMPT,
    TITLE_GEN_SYSTEM_PROMPT,
    FILENAME_SYSTEM_PROMPT
)
from .ppt.ppt_creator import PowerPointCreator
from .utils.image_caption_gen import ImageCaptionGenerator, find_image_shape, get_image_shape_from_path
from .utils.utils import format_simple_message_for_gpt, call_gpt_with_backoff, setup_logger


class PowerPointGenerator:
    def __init__(self, openai_key: str, model: str = "gpt-4"):
        openai.api_key = openai_key
        if model not in [model.id for model in openai.Model.list()['data']]:
            raise InvalidModel(
                f"Model {model} not found in list of models from OpenAI, make sure the "
                "name is correct or you have access to the model requested"
            )
        self.model = model
        self.image_gen = None
        self.logger = setup_logger(__name__)


    def create_powerpoint(self, user_input: str) -> str:
        """Generates a powerpoint based on the user's input"""
        # identify if the user passed in data for a chart
        data_messages = format_simple_message_for_gpt(
            system_message=CHART_DATA_IDENTIFICATION, user_message=user_input)
        data_response = call_gpt_with_backoff(
            messages=data_messages, model=self.model, temperature=0, max_length=MAX_CONTENT_LENGTH)

        # data was found in the input, determine which chart type fits the data best
        if data_response.lower() == "data found":
            best_chart_messages = format_simple_message_for_gpt(
                system_message=BEST_CHART_FOR_DATA_SYSTEM_PROMPT, user_message=user_input)
            best_chart_response = call_gpt_with_backoff(
                messages=best_chart_messages, model=self.model, temperature=0)
            # append this chart type instruction to the user input for deck creation
            user_input += f"\nUse chart type: {best_chart_response}"

        # create the deck based on the user input and load its json
        deck_messages = format_simple_message_for_gpt(
            system_message=DECK_CREATION_SYSTEM_PROMPT, user_message=user_input)
        deck_response = call_gpt_with_backoff(
            messages=deck_messages, model=self.model, temperature=0.2, max_length=MAX_CONTENT_LENGTH)
        deck_json = json.loads(deck_response)

        # create a fitting title for the deck
        title_messages = format_simple_message_for_gpt(
            system_message=TITLE_GEN_SYSTEM_PROMPT, user_message=deck_response)
        title_response = call_gpt_with_backoff(messages=title_messages, model=self.model)
        title_response = title_response.replace('"', '')

        # create a filename to save the deck as
        filename_message = format_simple_message_for_gpt(
            system_message=FILENAME_SYSTEM_PROMPT, user_message=title_response)
        filename_response = call_gpt_with_backoff(
            messages=filename_message, model=self.model, temperature=0)
        filename_response = filename_response.replace('"', '')

        # create the generated deck and save it to the specified filename
        ppt = PowerPointCreator(title=title_response, slides_content=deck_json)
        ppt.create(file_name=filename_response)
        return filename_response

    
    def _generate_caption(self, shape: BaseShape) -> str:
        """Generates caption for image"""
        try:
            image = Image.open(BytesIO(shape.image.blob))
            return self.image_gen.infer(image)
        except Exception as e:
            self.logger.error(f"Error generating caption for image: {e}")
            self.logger.exception(e)
            return ""


    def create_alt_text_for_powerpoint(self, ppt_path: str, device: str = "cuda"):
        """
        Generates alt text for a powerpoint, writes back to the original path
        Params:
        ::ppt_path: path to the powerpoint
        ::device: device to run on, options include ["cuda", "cpu", "mps"]
        """
        # load in ppt
        ppt = Presentation(ppt_path)

        # check if we need to instantiate a new caption gen
        if not self.image_gen:
            self.image_gen = ImageCaptionGenerator(device=device)

        # Find images and create a caption for them
        generate_captions = []
        for slide_index, slide in enumerate(ppt.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    caption = self._generate_caption(shape)
                    generate_captions.append({"slide": slide_index, "shape": shape_index, "caption": caption})
                    slide.shapes[shape_index]._element._nvXxPr.cNvPr.attrib['descr'] = caption
                elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                    if hasattr(shape, "image"):
                        caption = self._generate_caption(shape)
                        generate_captions.append({"slide": slide_index, "shape": shape_index, "caption": caption})
                        slide.shapes[shape_index]._element._nvXxPr.cNvPr.attrib['descr'] = caption
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    # if the shape is a group, check if there is a nested image
                    image_path = find_image_shape(shape=shape)
                    if image_path:
                        image_shape = get_image_shape_from_path(shape=shape, path=image_path)
                        caption = self._generate_caption(image_shape)
                        generate_captions.append({"slide": slide_index, "shape_path": image_path, "caption": caption})
                        image_shape._element._nvXxPr.cNvPr.attrib['descr'] = caption
        self.logger.info(f"Generated alt text: {generate_captions}")
        # save ppt
        ppt.save(ppt_path)
        self.logger.info("Updated Presentation alt-text")
