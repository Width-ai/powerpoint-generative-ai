from typing import List
from pptx import Presentation
from pptx.slide import Slide
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from powerpoint_generative_ai.utils.utils import setup_logger


SLIDE_LAYOUTS = {
    "Title Slide": 0,
    "Title and Content": 1,
    "Section Header": 2,
    "Two Content": 3,
    "Comparison": 4,
    "Title Only": 5,
    "Blank": 6,
    "Content with Caption": 7,
    "Picture with Caption": 8,
}


class PowerPointCreator:
    def __init__(self, title: str, slides_content: List[dict]):
        self.title = title
        self.slides_content = slides_content
        self.presentation = Presentation()
        self.logger = setup_logger(__name__)


    def create(self, file_name: str = 'presentation.pptx'):
        """Creates powerpoint and saves it to specified filename"""
        slide_layout = self.presentation.slide_layouts[0] 
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = self.title

        # Creating slides
        for content in self.slides_content:
            self.add_slide(content)

        self.save(file_name=file_name)


    def add_slide(self, content: dict):
        """Helper function to add slides to powerpoint"""
        text_content = content.get('content', None)
        chart_data = content.get('chart_data', None)
        image_path = content.get('diagram_name', None)
        table_data = content.get('table_data', None)

        LAYOUT = SLIDE_LAYOUTS['Title Slide']
        if table_data:
            LAYOUT = SLIDE_LAYOUTS['Title Only']
        elif text_content and (chart_data or image_path):
            LAYOUT = SLIDE_LAYOUTS['Two Content'] # Text column and blank right side
        elif chart_data and not text_content:
            LAYOUT = SLIDE_LAYOUTS['Title Only'] # Just a title for a big chart
        elif (text_content and not chart_data):
            LAYOUT = SLIDE_LAYOUTS['Title and Content'] # Standard slide format

        slide_layout = self.presentation.slide_layouts[LAYOUT]
        slide = self.presentation.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = content.get('title', None)

        
        if table_data:
            title.text = text_content
            title.text_frame.paragraphs[0].font.size = Pt(13)
            self.add_table(slide=slide, csv_text=table_data)
            return

        if text_content:
            content_shape = slide.shapes.placeholders[1]
            content_shape.text = text_content
        
        if chart_data:
            self.add_chart(data=chart_data, slide=slide, chart_type=content.get('chart_type', XL_CHART_TYPE.COLUMN_CLUSTERED))

        if image_path:
            self.add_image(image_path=image_path, slide=slide)



    def add_chart(self, data: dict, slide: Slide, x: Inches = Inches(4.75), y: Inches = Inches(2), cx: Inches=Inches(5.5), cy: Inches = Inches(4.5), chart_type: int = XL_CHART_TYPE.COLUMN_CLUSTERED):
        """Creates a chart and adds it to the current slide"""
        chart_data = CategoryChartData()
        chart_data.categories = data['categories']
        for series in data['series']:
            chart_data.add_series(series['name'], series['values'])

        chart = slide.shapes.add_chart(chart_type, x, y, cx, cy, chart_data).chart
        chart.has_legend = True

        for series in chart.series:
            series.has_data_labels = True

    def add_image(self, image_path: str, slide: Slide, x: Inches = Inches(4.75), y: Inches = Inches(2), cx: Inches=Inches(4), cy: Inches = Inches(3.5)):
        """Adds an image to the current slide"""
        slide.shapes.add_picture(image_path, x, y, cx, cy)

    def add_table(self, slide: Slide, csv_text: str):
        """Adds a table to the current slide"""

        shapes = slide.shapes
        self._csv_to_table(shapes, csv_text)

    def save(self, file_name: str):
        self.presentation.save(file_name)
        self.logger.info(f"Presentation successfully created: {file_name}")

    def _csv_to_table(self, shapes, csv_text):
        rows = csv_text.split('\n')
        data = [row.split(',') for row in rows]

        row_count = len(data)
        col_count = len(data[0])

        left = top = Inches(2.0)
        width = Inches(6.0)
        height = Inches(0.8)


        table = shapes.add_table(row_count, col_count, left, top, width, height).table

        for r in range(row_count):
            for c in range(col_count):
                cell = table.cell(r, c)
                cell.text = data[r][c]
